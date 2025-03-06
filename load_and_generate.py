import base64
import io
import os
import sys
from typing import *
from typing import List

import ollama
from langchain.text_splitter import RecursiveCharacterTextSplitter
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from tqdm import tqdm


def extract_image_from_shape(shape) -> Optional[Image.Image]:
    """
    Extracts image from a PowerPoint shape if it's an image.

    Args:
        shape: PowerPoint shape object

    Returns:
        PIL.Image or None: The extracted image if shape contains an image, None otherwise
    """
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image_blob = shape.image.blob
        return Image.open(io.BytesIO(image_blob))
    return None


def extract_content_from_slides(filepath: str) -> Tuple[str, List[Image.Image]]:
    """
    Extracts both text and images from a PowerPoint file.

    Args:
        filepath (str): The path to the PowerPoint file.

    Returns:
        Tuple[str, List[Image.Image]]: The extracted text and list of images from the slides.
    """
    prs = Presentation(filepath)
    text = ""
    images = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

            image = extract_image_from_shape(shape)
            if image:
                images.append(image)

    return text, images


def image_to_base64(image: Image.Image) -> str:
    """
    Convert a PIL Image to base64 string.

    Args:
        image (Image.Image): PIL Image object

    Returns:
        str: Base64 encoded string of the image
    """
    buffered = io.BytesIO()
    image.save(buffered, format="PNG")
    img_str = base64.b64encode(buffered.getvalue()).decode()
    return img_str


def process_images(images: List[Image.Image]) -> Dict[str, str]:
    """
    Process a list of PIL Image objects and return their descriptions.

    Args:
        images (List[Image.Image]): List of PIL Image objects to process.

    Returns:
        Dict[str, str]: Dictionary mapping image index to its description.
    """
    descriptions = {}
    for idx, image in tqdm(enumerate(images)):
        try:
            base64_image = image_to_base64(image)
            response = ollama.chat(
                model="llava",
                messages=[
                    {
                        "role": "system",
                        "content": """You are an advanced image analysis assistant designed for a software architect.
                                    Your task is to analyze images and provide precise descriptions in at most **three sentences**.
                                    
                                    **Description Rules:**
                                    - If the image is a **logo**, identify the company or brand name.
                                    - If the image is a **diagram**, explain its purpose and components.
                                    - If the image is a **chart or graph**, describe what it represents.
                                    - If the image is a **table**, summarize its key contents.
                                    - If the image is a **screenshot**, describe its interface or purpose.
                                    - If the image is a **map**, specify its location and details.
                                    - If the image is a **software icon**, describe what it represents.
                                    - If the image is a **cloud service icon**, specify the provider and service.
                                    
                                    **Response Rules:**
                                    - Limit your response to a **maximum of three sentences**.
                                    - If the image does not match any of the categories above, return an **empty string**.
                                    """,
                    },
                    {
                        "role": "user",
                        "content": "Describe this image:",
                        "images": [base64_image],
                    },
                ],
            )
            # Extract just the content from the response
            descriptions[f"image_{idx}"] = response["message"]["content"]
        except Exception as e:
            print(f"Error processing image {idx}: {str(e)}")
            descriptions[f"image_{idx}"] = ""

    return descriptions


def split_text(text: str) -> List[str]:
    """
    Split text into chunks of 1000 characters.

    Args:
        text (str): The text to split.

    Returns:
        List[str]: List of text chunks.
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000, chunk_overlap=200, length_function=len, add_start_index=True
    )
    return text_splitter.split_text(text=text)


def process_text_with_images(chunks: List, descriptions: Mapping) -> str:
    """
    Process text with image descriptions and return a final output.

    Args:
        text (str): The text to process.
        descriptions (Dict[str, str]): Dictionary mapping image index to its description.

    Returns:
        str: The final processed output.
    """
    all_text = "\n".join([chunk for chunk in chunks])
    all_image_descriptions = "\n".join(
        [description for description in descriptions.values()]
    )

    try:
        response = ollama.chat(
            model="llama3:latest",
            messages=[
                {
                    "role": "system",
                    "content": """You are a highly skilled assistant that extracts key insights from text and images.
                            Your task is to summarize the provided text and image descriptions.
                            
                            **Guidelines:**
                            - Identify the company name from the text, if mentioned.
                            - Generate a clear and concise summary of the text and image content.
                            - Follow the structured output format below:

                            **Output Format:**
                            Client: {Extracted company name from the text, if available, if not, mention high level. Example: Canadian Tech Company or American Health Company, or something else based on your understasnding}
                            Use case summary: {
                            Concise summary of the text and image descriptions, without mentioning the company name or the company description if you have mentioned it in the Company field. 
                            Focus on the impact and achievement and make the text concise and as short as possible without loosing meaningful output
                            Highlight any mention of the landing zone while maintaining the context, achievement, impact and contextual cohrency.
                            Make it as short as possible. Limit your self to 3 impactful sentences.
                            }
                            You will only output the Use case summary without any heading and without any extra unecessary annotations.
                        """,
                },
                {
                    "role": "user",
                    "content": f"""Summarize the following information:

                            **Text:** {all_text}

                            **Image Descriptions:** {all_image_descriptions}
                        """,
                },
            ],
        )

    except Exception as e:
        print(f"Error processing text: {str(e)}")
    return response["message"]


def has_title(slide) -> bool:
    """Check if slide has a title shape."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip().lower() in [
            "analysis summary",
            "analysis",
        ]:
            return True
    return False


def find_text_box(slide) -> Optional[Any]:
    """Find existing text box in slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            return shape
    return None


def append_or_create_text_box(
    slide, text: str, left: float, top: float, width: float, height: float
):
    """Append to existing text box or create new one."""
    text_box = find_text_box(slide)
    if text_box:
        # Append to existing text box
        paragraph = text_box.text_frame.add_paragraph()
        paragraph.text = text
    else:
        # Create new text box
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_box.text_frame.text = text
    return text_box


if __name__ == "__main__":
    # Check arguments
    if len(sys.argv) < 2:
        print("Usage: python script.py <input_pptx> [output_pptx]")
        exit(1)

    input_pptx = sys.argv[1]
    output_pptx = (
        sys.argv[2]
        if len(sys.argv) > 2
        else f"{os.path.splitext(input_pptx)[0]}_with_analysis.pptx"
    )

    text, images = extract_content_from_slides(input_pptx)
    print(f"\nExtracted {len(images)} images")

    # Save images and process them
    filename = os.path.basename(input_pptx).split(".")[0]
    for i, img in enumerate(images):
        img.save(f"{filename}_{i}.png")

    # Process images and print descriptions
    print("\nProcessing images...")
    descriptions = process_images(images)
    print("\nGenerating Image Descriptions:")
    for img_id, description in descriptions.items():
        print(description)

    print("\nSplitting text...")
    chunks = split_text(text)

    output = process_text_with_images(chunks, descriptions)["content"]

    with open(f"{filename}_analysis.txt", "a+") as f:
        f.write("\n")
        f.write(output)
